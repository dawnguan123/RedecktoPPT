#!/usr/bin/env python3
"""
RedecktoPPT - PPTX Renderer
组件化渲染引擎 v3.0

特性：
- 废除全页背景逻辑
- 组件化渲染：图片局部裁剪 + 文本独立框
- 语义化样式映射
- Z-Order 分层控制
"""

import os
from typing import List, Dict, Any, Optional
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PptxRenderer:
    """
    RedecktoPPT 组件化渲染引擎 v3.0
    
    核心设计：
    1. 不使用整页截图作为背景
    2. 图片：根据 bbox 从 PDF 实时裁剪局部图像
    3. 文本：根据 bbox 创建透明背景文本框
    4. 分层：图片在下层，文字在上层
    """
    
    # 布局索引
    LAYOUT_INDICES = {
        'cover': 6,
        'title': 0,
        'content': 1,
        'blank': 6,
    }
    
    # 样式配置
    STYLE_CONFIG = {
        'title': {
            'font_size': Pt(32),
            'bold': True,
            'color': RGBColor(0, 0, 0),
            'alignment': PP_ALIGN.LEFT
        },
        'text': {
            'font_size': Pt(14),
            'bold': False,
            'color': RGBColor(51, 51, 51),
            'alignment': PP_ALIGN.LEFT
        },
        'caption': {
            'font_size': Pt(10),
            'bold': False,
            'color': RGBColor(128, 128, 128),
            'alignment': PP_ALIGN.LEFT
        }
    }
    
    def __init__(
        self,
        template_path: Optional[str] = None,
        aspect_ratio: str = "16:9"
    ):
        # 初始化 PPT
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # 设置尺寸
        if aspect_ratio == "16:9":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
        elif aspect_ratio == "4:3":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        
        self.aspect_ratio = aspect_ratio
        self.slide_count = 0
        self._pdf_doc = None  # PDF 文档句柄
    
    def set_pdf_source(self, pdf_path: str):
        """设置 PDF 源，用于裁剪图片"""
        if os.path.exists(pdf_path):
            self._pdf_doc = fitz.open(pdf_path)
    
    def close(self):
        """关闭 PDF 文档"""
        if self._pdf_doc:
            self._pdf_doc.close()
    
    def create_slide(
        self,
        page_data: Dict,
        slide_index: int,
        total_pages: int,
        pdf_path: str = None
    ):
        """
        创建单页幻灯片（组件化渲染）
        
        Args:
            page_data: 页面数据，包含:
                - blocks: 文本块列表
                - images: 图片块列表
                - width: PDF 宽度 (points)
                - height: PDF 高度 (points)
            slide_index: 当前页码
            total_pages: 总页数
            pdf_path: PDF 文件路径（用于裁剪图片）
        """
        from core.coordinate_mapper import CoordinateMapper, SlideType
        
        # 1. 选择布局
        layout_idx = self._select_layout_index(slide_index, total_pages)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        
        # 2. 确定幻灯片类型
        if slide_index == 0:
            slide_type = SlideType.COVER
        elif slide_index == total_pages - 1:
            slide_type = SlideType.BACKCOVER
        else:
            slide_type = SlideType.NORMAL
        
        # 3. 获取页面数据
        blocks = page_data.get('blocks', [])
        images = page_data.get('images', [])
        pdf_width = page_data.get('width', 595)
        pdf_height = page_data.get('height', 842)
        
        # 4. 初始化坐标转换器
        margin = 0.3 if slide_type == SlideType.NORMAL else 0
        mapper = CoordinateMapper(
            pdf_width=pdf_width,
            pdf_height=pdf_height,
            ppt_width_inch=10 if self.aspect_ratio == "16:9" else 10,
            ppt_height_inch=5.625 if self.aspect_ratio == "16:9" else 7.5,
            margin_inch=margin
        )
        
        # 5. 打开 PDF（如果需要裁剪图片）
        if pdf_path and os.path.exists(pdf_path):
            if not self._pdf_doc:
                self._pdf_doc = fitz.open(pdf_path)
        
        # ========== 分层渲染 ==========
        
        # 第一层：渲染图片（底层）
        for img in images:
            self._render_image(
                slide=slide,
                img_data=img,
                mapper=mapper,
                page_index=slide_index,
                slide_type=slide_type
            )
        
        # 第二层：渲染文本（顶层）
        for block in blocks:
            self._render_text_block(
                slide=slide,
                block=block,
                mapper=mapper,
                slide_type=slide_type
            )
        
        self.slide_count += 1
        return slide
    
    def _render_image(
        self,
        slide,
        img_data: Dict,
        mapper,
        page_index: int,
        slide_type
    ):
        """
        渲染图片（局部裁剪）
        
        根据 bbox 从 PDF 页面中裁剪对应区域
        """
        try:
            bbox = img_data.get('bbox', [])
            if not bbox or len(bbox) != 4:
                return
            
            x1, y1, x2, y2 = bbox
            
            # 转换坐标
            left, top, width, height = mapper.to_pptx_geometry(
                x1, y1, x2 - x1, y2 - y1, slide_type
            )
            
            # 确保最小尺寸
            width = max(width, Emu(10000))
            height = max(height, Emu(10000))
            
            # 从 PDF 裁剪图片
            if self._pdf_doc and page_index < len(self._pdf_doc):
                page = self._pdf_doc[page_index]
                
                # 裁剪区域（PDF 坐标）
                clip_rect = fitz.Rect(x1, y1, x2, y2)
                
                # 渲染为图片
                zoom = 2.0  # 高清渲染
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, clip=clip_rect)
                
                # 保存到临时文件
                temp_path = f"/tmp/redeck_clip_{page_index}_{hash(str(bbox))}.png"
                pix.save(temp_path)
                
                # 添加到 PPT
                slide.shapes.add_picture(
                    temp_path,
                    left, top,
                    width=width,
                    height=height
                )
                
                # 清理临时文件
                try:
                    os.remove(temp_path)
                except:
                    pass
                    
        except Exception as e:
            print(f"   ⚠️ 图片裁剪失败: {e}")
    
    def _render_text_block(
        self,
        slide,
        block: Dict,
        mapper,
        slide_type
    ):
        """
        渲染文本块（独立文本框）
        
        根据 bbox 创建透明背景文本框
        """
        try:
            # 获取文本
            text = block.get('text', '').strip()
            if not text:
                return
            
            # 获取坐标
            bbox = block.get('bbox', [])
            if not bbox or len(bbox) != 4:
                return
            
            x1, y1, x2, y2 = bbox
            
            # 转换坐标
            left, top, width, height = mapper.to_pptx_geometry(
                x1, y1, x2 - x1, y2 - y1, slide_type
            )
            
            # 确保最小尺寸
            width = max(width, Emu(50000))
            height = max(height, Emu(20000))
            
            # 创建文本框
            txBox = slide.shapes.add_textbox(left, top, width, height)
            
            # 配置文本框
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # 设置文本
            p = tf.paragraphs[0]
            p.text = text
            
            # 应用样式
            block_type = block.get('type', 'text')
            self._apply_semantic_style(p, block_type, slide_type)
            
            # 设置文本框背景透明
            txBox.fill.background()
            
            # 设置边框透明
            line = txBox.line
            line.fill.background()
            line.color = RGBColor(255, 255, 255)
            
        except Exception as e:
            print(f"   ⚠️ 文本渲染失败: {e}")
    
    def _apply_semantic_style(
        self,
        paragraph,
        block_type: str,
        slide_type
    ):
        """
        应用语义化样式
        
        - title: 32pt, 加粗
        - text: 14pt, 标准
        - caption: 10pt, 灰色
        """
        # 获取样式配置
        if block_type == 'title':
            style = self.STYLE_CONFIG['title']
        elif block_type == 'caption':
            style = self.STYLE_CONFIG['caption']
        else:
            style = self.STYLE_CONFIG['text']
        
        # 应用字号
        paragraph.font.size = style['font_size']
        
        # 应用加粗
        paragraph.font.bold = style['bold']
        
        # 应用颜色
        paragraph.font.color.rgb = style['color']
        
        # 应用对齐
        paragraph.alignment = style['alignment']
    
    def _select_layout_index(self, index: int, total: int) -> int:
        """选择布局"""
        if index == 0:
            return self.LAYOUT_INDICES['cover']
        if index == total - 1:
            return self.LAYOUT_INDICES['cover']
        return self.LAYOUT_INDICES['content']
    
    def save(self, output_path: str):
        """保存文件"""
        # 关闭 PDF
        self.close()
        
        # 确保目录存在
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        
        self.prs.save(output_path)
        print(f"✨ RedecktoPPT 生成成功: {output_path}")
        print(f"   📊 共 {self.slide_count} 页")


def render_from_parsed_data(
    parsed_data: Dict,
    pdf_path: str,
    output_path: str,
    aspect_ratio: str = "16:9"
):
    """
    从解析数据渲染 PPT
    
    Args:
        parsed_data: MinerU 解析结果
        pdf_path: PDF 文件路径
        output_path: 输出路径
        aspect_ratio: 宽高比
    """
    renderer = PptxRenderer(aspect_ratio=aspect_ratio)
    renderer.set_pdf_source(pdf_path)
    
    pages = parsed_data.get('pages', [])
    
    for i, page in enumerate(pages):
        renderer.create_slide(
            page_data=page,
            slide_index=i,
            total_pages=len(pages),
            pdf_path=pdf_path
        )
    
    renderer.save(output_path)


if __name__ == "__main__":
    print("=" * 60)
    print("🧪 PptxRenderer v3.0 组件化渲染测试")
    print("=" * 60)
    
    # 测试数据
    test_page = {
        'width': 595,
        'height': 842,
        'blocks': [
            {
                'text': '智能协同：构建基于飞',
                'type': 'title',
                'bbox': [50, 50, 400, 80]
            },
            {
                'text': '这是一个测试段落，用于验证组件化渲染效果。',
                'type': 'text',
                'bbox': [50, 100, 500, 150]
            }
        ],
        'images': []
    }
    
    renderer = PptxRenderer()
    
    renderer.create_slide(
        page_data=test_page,
        slide_index=0,
        total_pages=1,
        pdf_path=None
    )
    
    renderer.save("test_component.pptx")
    print("\n✅ 测试完成!")
