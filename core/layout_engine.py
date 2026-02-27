"""
RedecktoPPT - Layout Engine
智能布局引擎：根据内容属性自动选择 PPT 布局

功能：
1. 布局类型枚举
2. 布局预测算法
3. PPTX 模板对接
4. 智能填充
"""

from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple
from enum import Enum


class SlideType(Enum):
    """幻灯片布局类型"""
    COVER = "cover"           # 封面
    POSTER = "poster"         # 海报页（单图+少量文字）
    CONTENT = "content"       # 内容页
    TABLE = "table"           # 表格页
    BACK_COVER = "backcover"  # 封底
    BLANK = "blank"           # 空白页


@dataclass
class LayoutConfig:
    """布局配置"""
    slide_type: SlideType
    title_text: str = ""
    body_text: str = ""
    images: List[Dict] = field(default_factory=list)
    tables: List[Dict] = field(default_factory=list)
    alignment: str = "left"  # left/center


@dataclass 
class ElementStats:
    """元素统计"""
    char_count: int = 0
    word_count: int = 0
    image_count: int = 0
    table_count: int = 0
    title_count: int = 0
    has_bullet: bool = False


class LayoutEngine:
    """
    布局引擎
    
    根据内容自动预测最佳布局
    """
    
    # 布局阈值配置
    POSTER_CHAR_THRESHOLD = 50   # 海报页字符上限
    POSTER_IMAGE_MIN = 1         # 海报页最少图片
    POSTER_IMAGE_MAX = 1          # 海报页最多图片
    
    # PPT Layout 索引（python-pptx）
    LAYOUT_INDICES = {
        SlideType.COVER: 6,      # blank (自定义)
        SlideType.POSTER: 6,     # blank (图片+标题)
        SlideType.CONTENT: 1,    # Title and Content
        SlideType.TABLE: 5,      # Title Only (手动插入表格)
        SlideType.BACK_COVER: 6, # blank (居中文字)
        SlideType.BLANK: 6,      # blank
    }
    
    def __init__(self):
        self.current_layout: Optional[SlideType] = None
    
    def predict_layout(
        self,
        elements: List[Dict[str, Any]],
        page_index: int = 0,
        total_pages: int = 1
    ) -> SlideType:
        """
        预测布局类型
        
        Args:
            elements: MinerU 元素列表
            page_index: 当前页码（从 0 开始）
            total_pages: 总页数
            
        Returns:
            SlideType: 预测的布局类型
        """
        # 统计元素
        stats = self._analyze_elements(elements)
        
        # 决策流程
        # 1. 封面（第一页）
        if page_index == 0:
            self.current_layout = SlideType.COVER
            return SlideType.COVER
        
        # 2. 封底（最后一页）
        if page_index == total_pages - 1:
            self.current_layout = SlideType.BACK_COVER
            return SlideType.BACK_COVER
        
        # 3. 表格页
        if stats.table_count > 0:
            self.current_layout = SlideType.TABLE
            return SlideType.TABLE
        
        # 4. 海报页（少量文字 + 1张图）
        if (stats.char_count < self.POSTER_CHAR_THRESHOLD and 
            stats.image_count == self.POSTER_IMAGE_MIN):
            self.current_layout = SlideType.POSTER
            return SlideType.POSTER
        
        # 5. 默认内容页
        self.current_layout = SlideType.CONTENT
        return SlideType.CONTENT
    
    def _analyze_elements(self, elements: List[Dict]) -> ElementStats:
        """分析元素统计"""
        stats = ElementStats()
        
        for elem in elements:
            elem_type = elem.get('type', '').lower()
            
            if elem_type == 'title':
                stats.title_count += 1
                stats.char_count += len(elem.get('content', ''))
                stats.word_count += len(elem.get('content', '').split())
            
            elif elem_type == 'text' or elem_type == 'body':
                content = elem.get('content', '')
                stats.char_count += len(content)
                stats.word_count += len(content.split())
                
                # 检测列表
                if any(c in content for c in ['•', '-', '*', '1.', '2.']):
                    stats.has_bullet = True
            
            elif elem_type == 'image' or elem_type == 'figure':
                stats.image_count += 1
            
            elif elem_type == 'table':
                stats.table_count += 1
        
        return stats
    
    def get_layout_index(self, slide_type: SlideType) -> int:
        """获取 PPT layout 索引"""
        return self.LAYOUT_INDICES.get(slide_type, 6)
    
    def create_layout_config(
        self,
        elements: List[Dict],
        layout_type: SlideType
    ) -> LayoutConfig:
        """
        创建布局配置
        
        Args:
            elements: 元素列表
            layout_type: 布局类型
            
        Returns:
            LayoutConfig: 填充配置
        """
        config = LayoutConfig(slide_type=layout_type)
        
        # 分类元素
        for elem in elements:
            elem_type = elem.get('type', '').lower()
            content = elem.get('content', '').strip()
            
            if elem_type in ('title', 'heading'):
                config.title_text = content
                
            elif elem_type in ('text', 'body', 'paragraph'):
                if config.body_text:
                    config.body_text += '\n' + content
                else:
                    config.body_text = content
                
                # 检测列表标记
                if content.startswith(('•', '-', '*')):
                    config.alignment = 'left'
            
            elif elem_type in ('image', 'figure'):
                config.images.append(elem)
            
            elif elem_type == 'table':
                config.tables.append(elem)
        
        return config


class PPTXLayoutFiller:
    """
    PPTX 布局填充器
    
    将内容填充到对应的 placeholder 中
    """
    
    def __init__(self, presentation):
        """
        初始化
        
        Args:
            presentation: python-pptx Presentation 对象
        """
        self.prs = presentation
    
    def add_slide(
        self,
        layout_type: SlideType,
        config: LayoutConfig
    ):
        """
        添加幻灯片并填充内容
        
        Args:
            layout_type: 布局类型
            config: 布局配置
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # 获取 layout 索引
        layout_idx = self._get_layout_index(layout_type)
        slide_layout = self.prs.slide_layouts[layout_idx]
        
        # 创建幻灯片
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 根据布局类型填充
        if layout_type == SlideType.COVER:
            self._fill_cover(slide, config)
        elif layout_type == SlideType.POSTER:
            self._fill_poster(slide, config)
        elif layout_type == SlideType.CONTENT:
            self._fill_content(slide, config)
        elif layout_type == SlideType.TABLE:
            self._fill_table(slide, config)
        elif layout_type == SlideType.BACK_COVER:
            self._fill_backcover(slide, config)
        
        return slide
    
    def _get_layout_index(self, layout_type: SlideType) -> int:
        """获取 layout 索引"""
        mapping = {
            SlideType.COVER: 6,
            SlideType.POSTER: 6,
            SlideType.CONTENT: 1,
            SlideType.TABLE: 5,
            SlideType.BACK_COVER: 6,
            SlideType.BLANK: 6,
        }
        return mapping.get(layout_type, 6)
    
    def _fill_cover(self, slide, config: LayoutConfig):
        """填充封面"""
        # 尝试获取 title placeholder
        title = self._get_placeholder(slide, 'title')
        if title:
            title.text = config.title_text or "封面标题"
        
        # 副标题
        subtitle = self._get_placeholder(slide, 'subtitle')
        if subtitle:
            subtitle.text = config.body_text or ""
    
    def _fill_poster(self, slide, config: LayoutConfig):
        """填充海报页（图片+居中标题）"""
        from pptx.util import Inches, Pt
        
        # 获取页面尺寸
        width = slide.shapes.width
        height = slide.shapes.height
        
        # 如果有图片，铺满整个页面
        if config.images:
            img = config.images[0]
            if 'path' in img:
                # 全屏图片
                slide.shapes.add_picture(
                    img['path'],
                    0, 0,
                    width=width,
                    height=height
                )
        
        # 居中标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            width - Inches(1), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = config.title_text or "海报标题"
        p.alignment = 1  # center
    
    def _fill_content(self, slide, config: LayoutConfig):
        """填充内容页"""
        # Title
        title = self._get_placeholder(slide, 'title')
        if title:
            title.text = config.title_text or "标题"
        
        # Content (body)
        content = self._get_placeholder(slide, 'content')
        if content:
            # 清除默认内容
            for p in content.text_frame.paragraphs:
                p.text = ""
            
            # 逐段填充
            lines = config.body_text.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = content.text_frame.paragraphs[0]
                else:
                    p = content.text_frame.add_paragraph()
                
                # 处理列表
                if line.strip().startswith(('•', '-', '*')):
                    p.level = 0
                    p.text = line.strip()[1:].strip()
                else:
                    p.text = line
    
    def _fill_table(self, slide, config: LayoutConfig):
        """填充表格页"""
        # Title
        title = self._get_placeholder(slide, 'title')
        if title:
            title.text = config.title_text or "表格"
        
        # 插入表格
        if config.tables:
            table_data = config.tables[0].get('data', [])
            if table_data:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 2
                
                # 插入表格 shape
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(3)
                
                table = slide.shapes.add_table(
                    rows, cols, left, top, width, height
                ).table
                
                # 填充数据
                for r, row in enumerate(table_data):
                    for c, cell in enumerate(row):
                        table.cell(r, c).text = str(cell)
    
    def _fill_backcover(self, slide, config: LayoutConfig):
        """填充封底"""
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        
        # 居中感谢语
        width = slide.shapes.width
        height = slide.shapes.height
        
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            width - Inches(2), Inches(1)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = config.body_text or "谢谢观看"
        p.alignment = PP_ALIGN.CENTER
    
    def _get_placeholder(self, slide, placeholder_type: str):
        """
        获取 placeholder
        
        Args:
            placeholder_type: 'title', 'content', 'subtitle', 'picture'
            
        Returns:
            Shape 或 None
        """
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if placeholder_type == 'title':
                if 'title' in str(ph_type).lower():
                    return shape
            elif placeholder_type == 'content':
                if 'body' in str(ph_type).lower():
                    return shape
            elif placeholder_type == 'subtitle':
                if 'subtitle' in str(ph_type).lower():
                    return shape
            elif placeholder_type == 'picture':
                if 'picture' in str(ph_type).lower():
                    return shape
        
        return None


# ========== 便捷函数 =========-

def predict_and_fill(
    presentation,
    elements: List[Dict],
    page_index: int = 0,
    total_pages: int = 1
):
    """
    一站式预测+填充
    
    Args:
        presentation: Presentation 对象
        elements: 元素列表
        page_index: 页码
        total_pages: 总页数
    """
    engine = LayoutEngine()
    filler = PPTXLayoutFiller(presentation)
    
    # 预测布局
    layout_type = engine.predict_layout(elements, page_index, total_pages)
    
    # 创建配置
    config = engine.create_layout_config(elements, layout_type)
    
    # 填充
    slide = filler.add_slide(layout_type, config)
    
    return slide


# ========== 单元测试 =========-

def test_layout_engine():
    """单元测试"""
    print("=" * 60)
    print("🧪 Layout Engine 单元测试")
    print("=" * 60)
    
    engine = LayoutEngine()
    
    # 测试 1: 封面检测
    print("\n📌 测试 1: 封面检测")
    elements = [
        {'type': 'title', 'content': '深度学习简介'},
        {'type': 'text', 'content': '第一章 基础知识'},
    ]
    layout = engine.predict_layout(elements, page_index=0, total_pages=5)
    print(f"   Page 0 → {layout.value}")
    assert layout == SlideType.COVER, "Should be COVER"
    print("   ✅ Pass")
    
    # 测试 2: 封底检测
    print("\n📌 测试 2: 封底检测")
    layout = engine.predict_layout(elements, page_index=4, total_pages=5)
    print(f"   Page 4 → {layout.value}")
    assert layout == SlideType.BACK_COVER, "Should be BACK_COVER"
    print("   ✅ Pass")
    
    # 测试 3: 表格页检测
    print("\n📌 测试 3: 表格页检测")
    elements = [
        {'type': 'title', 'content': '数据对比'},
        {'type': 'table', 'data': [['A', 'B'], [1, 2]]},
    ]
    layout = engine.predict_layout(elements, page_index=1, total_pages=5)
    print(f"   Has table → {layout.value}")
    assert layout == SlideType.TABLE, "Should be TABLE"
    print("   ✅ Pass")
    
    # 测试 4: 海报页检测
    print("\n📌 测试 4: 海报页检测")
    elements = [
        {'type': 'title', 'content': '架构图'},
        {'type': 'text', 'content': '神经网络结构'},
        {'type': 'image', 'path': 'diagram.png'},
    ]
    layout = engine.predict_layout(elements, page_index=1, total_pages=5)
    print(f"   <50 chars + 1 image → {layout.value}")
    assert layout == SlideType.POSTER, "Should be POSTER"
    print("   ✅ Pass")
    
    # 测试 5: 内容页检测
    print("\n📌 测试 5: 内容页检测")
    elements = [
        {'type': 'title', 'content': '机器学习概述'},
        {'type': 'text', 'content': '机器学习是人工智能的一个分支...' * 10},
    ]
    layout = engine.predict_layout(elements, page_index=1, total_pages=5)
    print(f"   Normal content → {layout.value}")
    assert layout == SlideType.CONTENT, "Should be CONTENT"
    print("   ✅ Pass")
    
    # 测试 6: 布局配置创建
    print("\n📌 测试 6: 布局配置创建")
    elements = [
        {'type': 'title', 'content': '测试标题'},
        {'type': 'text', 'content': '第一段内容'},
        {'type': 'text', 'content': '第二段内容'},
        {'type': 'image', 'path': 'img.png'},
    ]
    config = engine.create_layout_config(elements, SlideType.CONTENT)
    print(f"   Title: {config.title_text}")
    print(f"   Body: {config.body_text[:30]}...")
    print(f"   Images: {len(config.images)}")
    assert config.title_text == '测试标题'
    assert '第一段内容' in config.body_text
    print("   ✅ Pass")
    
    # 测试 7: 元素统计
    print("\n📌 测试 7: 元素统计")
    elements = [
        {'type': 'title', 'content': '标题'},
        {'type': 'text', 'content': '这是正文内容'},
        {'type': 'text', 'content': '• 列表项1'},
        {'type': 'image', 'path': 'a.png'},
        {'type': 'image', 'path': 'b.png'},
    ]
    stats = engine._analyze_elements(elements)
    print(f"   char_count: {stats.char_count}")
    print(f"   image_count: {stats.image_count}")
    print(f"   has_bullet: {stats.has_bullet}")
    assert stats.char_count > 0
    assert stats.image_count == 2
    assert stats.has_bullet == True
    print("   ✅ Pass")
    
    print("\n" + "=" * 60)
    print("🎉 所有测试通过!")
    print("=" * 60)


if __name__ == "__main__":
    test_layout_engine()
