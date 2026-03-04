#!/usr/bin/env python3
"""
RedecktoPPT - 简化版
从 PDF 提取图片，精确检测并覆盖文字/Logo

用法:
    python simple.py 输入.pdf 输出.pptx [底部高度]
"""

import sys
import os
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw
import numpy as np
import pytesseract
import cv2


def detect_text_boxes(img: Image.Image, bottom_height: int = 200) -> list:
    """OCR 检测文字 - 只检测底部水印区域"""
    width, height = img.size
    
    custom_config = r'--oem 3 --psm 6'
    data = pytesseract.image_to_data(
        img, 
        output_type=pytesseract.Output.DICT,
        config=custom_config
    )
    
    boxes = []
    n = len(data['text'])
    
    # 水印区域更靠下（留出缓冲区）
    watermark_start = height - bottom_height + 50
    
    for i in range(n):
        text = data['text'][i].strip()
        conf = float(data['conf'][i])
        
        if conf > 0:
            x = data['left'][i]
            y = data['top'][i]
            w = data['width'][i]
            h = data['height'][i]
            
            # 只覆盖明确在水印区域的文字
            if y > watermark_start:
                # 添加来源标记 'ocr'
                boxes.append((x, y, x + w, y + h, 'ocr'))
    
    return boxes


def detect_logo_precise(img: Image.Image, bottom_height: int = 200) -> list:
    """精确检测 Logo - 针对右下角区域"""
    width, height = img.size
    boxes = []
    
    y_start = height - bottom_height
    
    # 右下角区域 - 进一步扩大检测范围
    corner_w = min(500, int(width * 0.45))
    corner_h = min(150, bottom_height + 50)
    
    corner_region = np.array(img.crop((
        width - corner_w, 
        height - corner_h, 
        width, 
        height
    )))
    
    gray_corner = cv2.cvtColor(corner_region, cv2.COLOR_RGB2GRAY)
    
    # 边缘检测
    canny = cv2.Canny(gray_corner, 15, 60)
    contours, _ = cv2.findContours(canny, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    for cnt in contours:
        area = cv2.contourArea(cnt)
        if area > 30:  # 更低的阈值
            x, y, w, h = cv2.boundingRect(cnt)
            # 只保留底部附近的
            box_y1 = height - corner_h + y + h
            if box_y1 > height - 50:
                boxes.append((
                    width - corner_w + x,
                    height - corner_h + y,
                    width - corner_w + x + w,
                    height - corner_h + y + h,
                    'logo'
                ))
    
    # 如果检测太少，扩大区域重新检测
    if len(boxes) < 2:
        # 尝试更宽的区域
        larger_w = min(400, width // 2)
        larger_h = min(150, bottom_height + 50)
        
        larger_region = np.array(img.crop((
            width - larger_w, 
            height - larger_h, 
            width, 
            height
        )))
        
        gray_large = cv2.cvtColor(larger_region, cv2.COLOR_RGB2GRAY)
        canny_large = cv2.Canny(gray_large, 15, 60)
        contours_large, _ = cv2.findContours(canny_large, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        for cnt in contours_large:
            area = cv2.contourArea(cnt)
            if area > 30:
                x, y, w, h = cv2.boundingRect(cnt)
                # 只保留右下部分 且 在底部附近
                if x > larger_w * 0.3:
                    box_y1 = height - larger_h + y + h
                    if box_y1 > height - 50:
                        boxes.append((
                            width - larger_w + x,
                            height - larger_h + y,
                            width - larger_w + x + w,
                            height - larger_h + y + h,
                            'logo'
                        ))
    
    return boxes


def merge_nearby_boxes(boxes: list, distance_threshold: int = 20, width: int = 0) -> list:
    """
    合并 box
    - Logo 盒子（只有 logo 来源）且在右下角：全部合并
    - 其他：按距离合并
    """
    if not boxes:
        return []
    
    boxes = sorted(boxes, key=lambda b: b[0])
    
    # 分类：纯 Logo vs 其他
    pure_logo_boxes = [b for b in boxes if b[4] == 'logo']
    other_boxes = [b for b in boxes if b[4] != 'logo']
    
    # 处理纯 Logo 盒子：如果都在右下角，合并为一个
    if pure_logo_boxes and width > 0:
        # 检查是否都在右下角（x > 60% width）
        right_boxes = [b for b in pure_logo_boxes if b[0] > width * 0.6]
        if len(right_boxes) >= len(pure_logo_boxes) * 0.7:  # 70% 以上在右侧
            # 全部合并，并扩展到边界
            x0 = min(b[0] for b in pure_logo_boxes)
            y0 = min(b[1] for b in pure_logo_boxes)
            x1 = max(b[2] for b in pure_logo_boxes)
            y1 = max(b[3] for b in pure_logo_boxes)
            
            # 如果右边还有空白，扩展到边界
            if width - x1 < 100:
                x1 = width
            
            merged = [(x0, y0, x1, y1, {'logo'})]
        else:
            merged = []
    else:
        merged = []
    
    # 处理其他盒子（按距离合并）
    if other_boxes:
        current_group = [other_boxes[0]]
        for i in range(1, len(other_boxes)):
            box = other_boxes[i]
            prev = current_group[-1]
            
            if box[0] - prev[2] < distance_threshold:
                current_group.append(box)
            else:
                x0 = min(b[0] for b in current_group)
                y0 = min(b[1] for b in current_group)
                x1 = max(b[2] for b in current_group)
                y1 = max(b[3] for b in current_group)
                sources = set(b[4] for b in current_group)
                merged.append((x0, y0, x1, y1, sources))
                current_group = [box]
        
        if current_group:
            x0 = min(b[0] for b in current_group)
            y0 = min(b[1] for b in current_group)
            x1 = max(b[2] for b in current_group)
            y1 = max(b[3] for b in current_group)
            sources = set(b[4] for b in current_group)
            merged.append((x0, y0, x1, y1, sources))
    
    return merged


def expand_boxes_to_cover_line(boxes: list, width: int, height: int, bottom_height: int = 200) -> list:
    """
    扩展策略：
    - 只有纯 OCR 检测（无 Logo）才扩展
    - 如果合并了 Logo 来源，不扩展
    """
    expanded = []
    
    # 底部区域边界
    bottom_start = height - bottom_height
    watermark_zone = height - 30  # 水印区域
    
    for box in boxes:
        x0, y0, x1, y1, sources = box
        box_width = x1 - x0
        
        # 只在底部区域考虑扩展
        if y0 > bottom_start:
            has_ocr = 'ocr' in sources
            has_logo = 'logo' in sources
            in_watermark_zone = y0 > watermark_zone
            
            # 只有纯 OCR（无 Logo）且在水印区域才扩展
            if has_ocr and not has_logo and in_watermark_zone and box_width < width * 0.25:
                expanded.append((0, y0, width, y1, sources))
            else:
                expanded.append(box)
        else:
            expanded.append(box)
    
    return expanded


def fill_boxes(img: Image.Image, boxes: list) -> None:
    """用周围颜色填充 - 从 Logo 下方采样"""
    draw = ImageDraw.Draw(img)
    width, height = img.size
    
    for box in boxes:
        x0, y0, x1, y1, sources = box
        
        # 扩展边距
        margin = 5
        x0 = max(0, x0 - margin)
        y0 = max(0, y0 - margin)
        x1 = min(width, x1 + margin)
        y1 = min(height, y1 + margin)
        
        # 优先从 Logo 下方采样
        samples = []
        
        # 下方采样（Logo 在底部，下方更可靠）
        if y1 < height - 10:
            for dx in range(x0, min(x1, width), max(1, (x1-x0)//5)):
                if y1 + 5 < height:
                    samples.append(img.getpixel((dx, y1 + 5)))
        
        # 如果下方采样不够，用上方
        if len(samples) < 3 and y0 > 10:
            for dx in range(x0, min(x1, width), max(1, (x1-x0)//5)):
                samples.append(img.getpixel((dx, y0 - 5)))
        
        if samples:
            avg = np.mean(samples, axis=0)
            fill = tuple(avg.astype(int))
        else:
            fill = (255, 255, 255)
        
        draw.rectangle([x0, y0, x1, y1], fill=fill)


def process_page(pdf_path: str, page_num: int, dpi: int = 150, bottom_height: int = 200) -> str:
    """处理单个页面"""
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat)
    
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    width, height = img.size
    
    all_boxes = []
    
    # OCR
    print(f"      OCR...")
    text_boxes = detect_text_boxes(img, bottom_height)
    print(f"      OCR: {len(text_boxes)} 个")
    all_boxes.extend(text_boxes)
    
    # Logo 检测
    print(f"      Logo检测...")
    logo_boxes = detect_logo_precise(img, bottom_height)
    print(f"      Logo: {len(logo_boxes)} 个")
    all_boxes.extend(logo_boxes)
    
    # 合并
    boxes = merge_nearby_boxes(all_boxes, distance_threshold=20, width=width)
    print(f"      合并: {len(boxes)} 个")
    
    # 扩展小区域
    boxes = expand_boxes_to_cover_line(boxes, width, height, bottom_height)
    print(f"      扩展后: {len(boxes)} 个")
    
    if boxes:
        fill_boxes(img, boxes)
    
    img_path = f"/tmp/redeck_page_{page_num}.png"
    img.save(img_path)
    
    doc.close()
    return img_path


def pdf_to_images(pdf_path: str, dpi: int = 150, bottom_height: int = 200) -> list:
    images = []
    doc = fitz.open(pdf_path)
    num_pages = len(doc)
    doc.close()
    
    for page_num in range(num_pages):
        img_path = process_page(pdf_path, page_num, dpi, bottom_height)
        images.append(img_path)
        print(f"   📄 处理第 {page_num + 1}/{num_pages} 页")
    
    return images


def create_ppt(images: list, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for img_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            width=Inches(10), height=Inches(5.625)
        )
        print(f"   🖼️ 添加: {Path(img_path).name}")
    
    prs.save(output_path)


def cleanup(images: list):
    for img in images:
        try:
            os.remove(img)
        except:
            pass


def main():
    pdf_path = "/Users/guanliming/Downloads/Data_Synergy.pdf"
    output_path = "/Users/guanliming/.openclaw/workspace/RedecktoPPT/output/Data_Synergy.pptx"
    bottom_height = 200
    
    if len(sys.argv) >= 2:
        pdf_path = sys.argv[1]
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    if len(sys.argv) >= 4:
        bottom_height = int(sys.argv[3])
    
    if not os.path.exists(pdf_path):
        print(f"❌ 文件不存在: {pdf_path}")
        sys.exit(1)
    
    # PPTX 转 PDF
    is_pptx = pdf_path.lower().endswith('.pptx')
    if is_pptx:
        print(f"🔄 检测到 PPTX，先转换为 PDF...")
        temp_pdf = f"/tmp/converted_{os.path.basename(pdf_path)}.pdf"
        
        # 调用 pptx_to_pdf.py
        import subprocess
        script_dir = os.path.dirname(os.path.abspath(__file__))
        result = subprocess.run(
            [sys.executable, os.path.join(script_dir, 'pptx_to_pdf.py'), 
             pdf_path, temp_pdf],
            capture_output=True, text=True
        )
        
        if result.returncode != 0 or not os.path.exists(temp_pdf):
            print(f"❌ PPTX 转换失败: {result.stderr}")
            sys.exit(1)
        
        pdf_path = temp_pdf
        print(f"   ✅ PPTX 转换为 PDF 成功")
    
    print(f"🎬 开始转换: {Path(pdf_path).name}")
    print(f"   底部检测高度: {bottom_height}px")
    print("-" * 40)
    
    print("📄 提取页面...")
    images = pdf_to_images(pdf_path, bottom_height=bottom_height)
    print(f"   ✅ 处理了 {len(images)} 页")
    
    print("\n🎨 生成 PPT...")
    create_ppt(images, output_path)
    
    cleanup(images)
    print("\n🎉 完成!")


if __name__ == "__main__":
    main()
